from io import BytesIO
from pathlib import Path
import re

import cairosvg
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


workspace = Path(__file__).parent
html = (workspace / "yankee-park-map.html").read_text(encoding="utf-8")
svg_match = re.search(r"<svg\b.*?</svg>", html, flags=re.DOTALL)
if svg_match is None:
    raise RuntimeError("The source HTML does not contain an SVG map.")

svg = svg_match.group(0)
map_png = cairosvg.svg2png(bytestring=svg.encode("utf-8"), output_width=1700, output_height=2200)
presentation = Presentation()
presentation.slide_width = Inches(8.5)
presentation.slide_height = Inches(11)
slide = presentation.slides.add_slide(presentation.slide_layouts[6])

background = slide.background.fill
background.solid()
background.fore_color.rgb = RGBColor(244, 246, 249)

title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(7.4), Inches(0.48))
title_frame = title_box.text_frame
title_frame.clear()
title = title_frame.paragraphs[0]
title.text = "Yankee Park - Lot Mapping"
title.font.name = "Aptos Display"
title.font.size = Pt(22)
title.font.bold = True
title.font.color.rgb = RGBColor(44, 62, 80)
title.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.72), Inches(7.4), Inches(0.25))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.clear()
subtitle = subtitle_frame.paragraphs[0]
subtitle.text = "Interactive SVG map converted from the source HTML"
subtitle.font.name = "Aptos"
subtitle.font.size = Pt(9)
subtitle.font.color.rgb = RGBColor(100, 116, 139)
subtitle.alignment = PP_ALIGN.CENTER

map_width = Inches(7.25)
map_height = Inches(9.39)
map_left = (presentation.slide_width - map_width) // 2
map_top = Inches(1.1)
map_shape = slide.shapes.add_picture(BytesIO(map_png), map_left, map_top, width=map_width, height=map_height)
if map_shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
    raise RuntimeError("The SVG map was not embedded as a picture.")

output = workspace / "yankee-park-map.pptx"
presentation.save(output)
print(output)